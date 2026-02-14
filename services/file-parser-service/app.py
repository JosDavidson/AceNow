from flask import Flask, jsonify, request
from flask_cors import CORS
import pdfplumber
from pptx import Presentation
import os
import re
from dotenv import load_dotenv

load_dotenv()

app = Flask(__name__)
CORS(app)

def clean_extracted_text(text):
    """Remove common PDF/binary junk and normalize whitespace"""
    if not text:
        return ""
    
    # Remove obvious PDF structural markers if they leaked
    text = re.sub(r'endstream|endobj|\d+ \d+ obj|<<|>>|stream', '', text)
    
    # Remove non-printable characters except common punctuation/newlines
    # Keep standard ASCII and some useful Unicode
    text = "".join(char for char in text if char.isprintable() or char in "\n\t\r")
    
    # Normalize whitespace
    text = re.sub(r'\s+', ' ', text).strip()
    return text

@app.route('/health', methods=['GET'])
def health():
    """Health check endpoint"""
    return jsonify({"status": "healthy", "service": "file-parser-service"}), 200

@app.route('/parse-file', methods=['POST'])
def parse_file():
    """Parse uploaded file and extract text"""
    try:
        if 'file' not in request.files:
            return jsonify({"success": False, "error": "No file uploaded"}), 400

        file = request.files['file']
        filename = file.filename.lower()
        extracted_text = ""

        # Parse PDF files using pdfplumber (more robust)
        if filename.endswith('.pdf'):
            try:
                with pdfplumber.open(file) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            extracted_text += page_text + "\n"
            except Exception as pdf_err:
                print(f"pdfplumber failed: {pdf_err}")
                return jsonify({"success": False, "error": f"PDF parsing failed: {str(pdf_err)}"}), 500

        # Parse PPTX files
        elif filename.endswith('.pptx'):
            try:
                prs = Presentation(file)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            extracted_text += shape.text + "\n"
            except Exception as ppt_err:
                print(f"PPTX parsing failed: {ppt_err}")
                return jsonify({"success": False, "error": f"PPTX parsing failed: {str(ppt_err)}"}), 500

        # Parse text files
        else:
            try:
                extracted_text = file.read().decode("utf-8", errors="ignore")
            except Exception as txt_err:
                return jsonify({"success": False, "error": f"Text file parsing failed: {str(txt_err)}"}), 500

        # Clean the text to ensure AI doesn't get junk
        cleaned_text = clean_extracted_text(extracted_text)
        
        print(f"Parsed file '{filename}' (Original: {len(extracted_text)}, Cleaned: {len(cleaned_text)})")
        
        if not cleaned_text or len(cleaned_text) < 10:
             return jsonify({
                "success": False, 
                "error": "No readable text found in file. Please ensure the file is not just images."
            }), 400

        return jsonify({
            "success": True,
            "text": cleaned_text,
            "filename": filename,
            "length": len(cleaned_text)
        }), 200

    except Exception as e:
        print(f"Parse error: {str(e)}")
        return jsonify({"success": False, "error": str(e)}), 500

if __name__ == '__main__':
    port = int(os.getenv('PORT', 5002))
    print(f"File Parser Service running on port {port}")
    app.run(host='0.0.0.0', port=port, debug=True)
